#!/usr/bin/env python
"""
Phase 9: Create Presentation Deck
Generate PowerPoint presentation for stakeholder presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_title_slide(prs, title, subtitle):
    """Create a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 55, 100)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

    return slide


def create_content_slide(prs, title, content_points):
    """Create a standard content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 55, 100)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(12)
        p.space_after = Pt(6)
        p.level = 0

    return slide


def create_metrics_slide(prs):
    """Create model comparison metrics slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Model Performance Comparison"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 55, 100)

    # Table
    rows, cols = 4, 6
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(3.5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Header
    headers = ["Model", "AUC-ROC", "AUC-PR", "Precision", "Recall", "F1 Score"]
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(25, 55, 100)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True

    # Data rows
    data = [
        ["Heuristic Rule", "N/A", "N/A", "80.3%", "72.8%", "0.764"],
        ["Logistic Regression", "0.999", "1.000", "99.5%", "99.0%", "0.993"],
        ["XGBoost", "0.999", "0.999", "98.1%", "100.0%", "0.990"],
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = cell_text
            if row_idx == 2:  # Highlight logistic regression
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(200, 230, 255)

    return slide


def create_recommendation_slide(prs):
    """Create deployment recommendation slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Deployment Recommendation"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 55, 100)

    # Phase 1
    phase1_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(3))
    phase1_frame = phase1_box.text_frame
    phase1_frame.word_wrap = True

    p = phase1_frame.paragraphs[0]
    p.text = "PHASE 1: Logistic Regression"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 55, 100)

    for point in [
        "Launch immediately",
        "F1: 0.993 (excellent)",
        "99.5% precision",
        "Fast (<1ms inference)",
        "Interpretable",
        "GBP 300-450K year 1"
    ]:
        p = phase1_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.level = 1
        p.space_before = Pt(6)

    # Phase 2
    phase2_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(3))
    phase2_frame = phase2_box.text_frame
    phase2_frame.word_wrap = True

    p = phase2_frame.paragraphs[0]
    p.text = "PHASE 2: XGBoost (Optional)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 55, 100)

    for point in [
        "Upgrade after Phase 1 validation",
        "100% recall (no false negatives)",
        "SHAP explanations",
        "A/B test for 3 months",
        "Upgrade if 5% uplift confirmed",
        "Additional GBP 3-5K saving"
    ]:
        p = phase2_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.level = 1
        p.space_before = Pt(6)

    return slide


def main():
    """Create full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating presentation...")

    # Slide 1: Title
    create_title_slide(
        prs,
        "FreshBox Churn Prediction",
        "Machine Learning Model for Customer Retention"
    )

    # Slide 2: Problem Statement
    create_content_slide(
        prs,
        "The Problem",
        [
            "Current State: 67.4% of customers churn within 6 months",
            "Impact: Major revenue loss from preventable churn",
            "Opportunity: Proactive retention targeting",
            "Goal: Build predictive model to identify at-risk customers",
            "Target: Enable retention team to act on top 5,000 at-risk customers"
        ]
    )

    # Slide 3: Data & Approach
    create_content_slide(
        prs,
        "Data & Methodology",
        [
            "Dataset: 1,500 customers with 49 features",
            "Features: Tenure, engagement, momentum, economic signals, friction, demographics",
            "Approach: Build 3 models (heuristic baseline, logistic, XGBoost)",
            "Validation: Stratified 80/20 train/test split with data leakage prevention",
            "Key Metric: Balance precision (budget efficiency) and recall (prevent churn)"
        ]
    )

    # Slide 4: Model Comparison
    create_metrics_slide(prs)

    # Slide 5: Heuristic Baseline
    create_content_slide(
        prs,
        "Baseline: Heuristic Rule",
        [
            "Simple 3-signal rule: Low rating OR declining frequency OR complaints",
            "Performance: F1 0.764, Precision 80.3%, Recall 72.8%",
            "Problem: Misses 55 churners (27.2% false negatives)",
            "Cost: Retention team wastes effort on 36 false positives",
            "Verdict: ML models can significantly improve upon this baseline"
        ]
    )

    # Slide 6: Logistic Regression
    create_content_slide(
        prs,
        "Model 1: Logistic Regression",
        [
            "Interpretable linear model: Understand each feature's impact on churn",
            "Top 5 Protective Factors: total_orders, order_completion_rate, recipes_rated, frequency_trend, late_order_freq",
            "Top 2 Risk Factors: support_tickets, skips",
            "Performance: F1 0.993, AUC-ROC 0.999, Precision 99.5%, Recall 99.0%",
            "Advantage: Catches 200/202 churners with only 1 false positive"
        ]
    )

    # Slide 7: XGBoost
    create_content_slide(
        prs,
        "Model 2: XGBoost (Non-Linear)",
        [
            "Advanced gradient boosting: Captures non-linear relationships",
            "Top Feature: late_order_freq (0.59 importance) - recent momentum is critical",
            "SHAP Explanations: Understand individual prediction drivers",
            "Performance: F1 0.990, AUC-ROC 0.999, Precision 98.1%, Recall 100.0%",
            "Advantage: Perfect recall - catches ALL 202 churners (0 false negatives)"
        ]
    )

    # Slide 8: Recommendation
    create_recommendation_slide(prs)

    # Slide 9: Business Impact
    create_content_slide(
        prs,
        "Expected Business Impact",
        [
            "Year 1 Revenue Protection: GBP 300-450K from churn prevention",
            "Phase 1 (Logistic): Identify ~1,484 likely churners from 5,000 flagged",
            "Retention Focus: Target declining frequency customers early (strongest signal)",
            "Phase 2 Upside: Additional GBP 3-5K if XGBoost upgrade justified by A/B testing",
            "Implementation: Batch scoring of 80,000 customers weekly via existing analytics pipeline"
        ]
    )

    # Slide 10: Next Steps
    create_content_slide(
        prs,
        "Implementation Timeline",
        [
            "Week 1-2: Deploy logistic regression scoring in production",
            "Week 3-4: Train retention team on top features and targeting strategies",
            "Month 1-3: Monitor AUC-ROC, precision, and churn rate impact weekly",
            "Month 2-3: A/B test XGBoost vs Logistic in parallel (if retention team capacity allows)",
            "Month 4: Review results and decide on Phase 2 (XGBoost) upgrade"
        ]
    )

    # Save
    output_path = "deck/FreshBox_Churn_Prediction.pptx"
    prs.save(output_path)
    print(f"[OK] Presentation saved to {output_path}")


if __name__ == "__main__":
    main()
